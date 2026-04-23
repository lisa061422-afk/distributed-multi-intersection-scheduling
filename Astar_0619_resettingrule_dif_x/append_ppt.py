"""
append_ppt.py — appends Simulation Test 2 (50 trials) slides to existing PPT.
Does NOT touch or regenerate existing slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colours (same light theme as make_ppt.py) ────────────────────────────────
C_BLUE   = RGBColor(0x1F, 0x6B, 0xB0)
C_GREEN  = RGBColor(0x1E, 0x8A, 0x44)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_YELLOW = RGBColor(0xD3, 0x7A, 0x00)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0x2C, 0x3E, 0x50)
C_MGRAY  = RGBColor(0x55, 0x6B, 0x82)

PPTX_PATH = r"C:\Users\robin\OneDrive\Documents\Github_file\Astar_0619_resettingrule_dif_x\WeakRule_Comparison.pptx"

prs = Presentation(PPTX_PATH)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────────
def bg(slide, color=RGBColor(0xFF,0xFF,0xFF)):
    sh = slide.shapes.add_shape(1, 0, 0, W, H)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()

def rect(slide, l, t, w, h, fill, line=None, line_w=Pt(0)):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line; sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def box(slide, l, t, w, h, text, font_size=14, bold=False,
        fg=C_LGRAY, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(font_size); run.font.bold = bold
    run.font.color.rgb = fg
    return txb

def title_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, C_BLUE)
    box(slide, 0.3, 0.1, 12, 0.6, title,
        font_size=28, bold=True, fg=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        box(slide, 0.3, 0.65, 12, 0.4, subtitle,
            font_size=14, fg=RGBColor(0xD0,0xE8,0xFF), align=PP_ALIGN.LEFT)

# ── data ─────────────────────────────────────────────────────────────────────
rows_all = [
    (1,  5,"[4 3 7 9 6]",   "2.94/2.05/1.44/1.18/1.03","0.5190","0.5190","YES", 25,   73,  2,  10, 0.27, 0.17,""),
    (2,  5,"[6 1 5 9 3]",   "0.53/1.59/1.60/1.90/2.55","1.3554","1.3554","YES", 36,   73,  5,  18, 0.10, 0.13,""),
    (3,  5,"[8 9 4 5 3]",   "0.88/1.89/0.28/1.30/1.29","2.7792","2.9938","NO★", 52,  279,  9,  61, 0.08, 0.17,""),
    (4,  4,"[6 4 6 11]",    "2.83/1.51/1.87/0.35",     "0.6629","0.6629","YES", 24,   48,  4,  11, 0.02, 0.02,""),
    (5,  4,"[5 11 4 6]",    "2.96/1.56/1.84/0.36",     "0.0000","0.0000","YES", 13,   13,  1,   1, 0.00, 0.00,""),
    (6,  6,"[8 7 5 4 6 9]", "2.63/1.53/2.01/1.76/1.87/2.02","NaN","NaN","N/A",4182,6357,  0,   0,30.01,30.00,"ON OFF"),
    (7,  6,"[1 10 3 3 7 2]","2.66/1.88/2.17/0.05/1.78/1.67","4.3281","NaN","N/A",5253,4721,322,0,26.41,30.00,"OFF"),
    (8,  3,"[2 9 4]",       "2.08/1.66/1.17",          "0.0000","0.0000","YES", 10,   10,  1,   1, 0.00, 0.00,""),
    (9,  6,"[11 5 1 4 5 9]","2.99/1.07/2.29/1.78/2.08/0.45","3.0768","NaN","N/A",407,5096, 49,  0, 0.40,30.02,"OFF"),
    (10, 4,"[3 5 7 8]",     "0.32/0.39/0.97/1.98",     "1.1253","1.1253","YES", 44,   73,  7,  16, 0.06, 0.06,""),
    (11, 6,"[7 11 5 4 5 3]","2.49/1.02/1.66/1.74/1.56/0.01","2.8541","NaN","N/A",385,4188, 44,  0, 0.44,30.00,"OFF"),
    (12, 6,"[11 3 4 7 11 12]","0.77/1.69/2.42/1.18/2.19/0.48","3.0143","NaN","N/A",1096,4967,197, 0, 3.03,30.00,"OFF"),
    (13, 5,"[11 12 1 6 3]", "1.35/1.64/0.28/0.89/2.78","1.2162","1.2162","YES", 28,  134,  4,  31, 0.02, 0.05,""),
    (14, 5,"[6 10 9 1 9]",  "2.52/0.50/2.34/0.86/0.92","1.0343","1.0343","YES", 37,  230,  6,  57, 0.03, 0.12,""),
    (15, 5,"[2 8 11 9 6]",  "1.31/2.30/1.70/0.25/1.75","0.2188","0.2188","YES", 20,   47,  2,   8, 0.01, 0.02,""),
    (16, 6,"[5 12 10 7 10 1]","2.58/2.46/2.73/0.39/0.25/0.42","2.5894","NaN","N/A",1928,6142,224, 0, 3.05,30.02,"OFF"),
    (17, 4,"[6 7 2 3]",     "2.43/1.40/2.42/0.02",     "0.4491","0.4491","YES", 16,   22,  2,   4, 0.02, 0.02,""),
    (18, 5,"[12 7 3 9 5]",  "2.01/0.09/1.91/0.10/2.23","1.0343","1.0343","YES", 36,  140,  4,  27, 0.04, 0.11,""),
    (19, 4,"[2 7 1 8]",     "2.99/2.31/1.72/0.31",     "0.1985","0.1985","YES", 40,   46,  5,   7, 0.03, 0.01,""),
    (20, 5,"[8 1 10 7 6]",  "2.36/1.23/1.44/0.54/0.96","1.6664","1.6664","YES",322,  746, 32, 142, 0.28, 0.90,""),
    (21, 6,"[3 6 12 3 12 12]","0.27/1.39/1.51/0.94/0.14/0.73","0.5268","0.5268","YES",112,4295, 22,1535, 0.14,26.31,""),
    (22, 3,"[3 10 11]",     "0.13/0.91/2.94",          "0.0000","0.0000","YES", 10,   10,  1,   1, 0.00, 0.00,""),
    (23, 5,"[8 1 6 12 5]",  "0.29/1.39/2.89/1.03/2.40","1.7178","1.7974","NO★", 43, 2584,  6, 656, 0.02, 5.92,""),
    (24, 6,"[3 6 9 5 3 12]","1.95/2.60/0.08/0.80/1.51/0.20","0.3394","0.3394","YES", 19,  27,  2,   5, 0.01, 0.01,""),
    (25, 6,"[3 5 3 2 3 4]", "1.90/0.84/1.09/0.02/1.10/1.60","3.9051","NaN","N/A",3990,6579,691,  0,11.79,30.00,"OFF"),
    (26, 3,"[8 4 8]",       "0.08/2.66/0.05",          "0.9698","0.9698","YES", 17,   17,  2,   2, 0.01, 0.00,""),
    (27, 3,"[10 1 9]",      "2.91/2.62/2.13",          "0.4873","0.4873","YES", 13,   13,  2,   2, 0.01, 0.00,""),
    (28, 6,"[6 11 5 12 2 12]","2.50/2.54/0.37/1.79/0.05/2.16","1.5560","NaN","N/A",250,5757, 42,  0, 0.11,30.01,"OFF"),
    (29, 3,"[2 3 11]",      "1.09/1.62/1.70",          "0.6427","0.6427","YES", 20,   35,  4,   9, 0.00, 0.01,""),
    (30, 3,"[7 8 4]",       "1.26/1.36/2.80",          "0.6820","0.6820","YES", 23,   23,  3,   3, 0.01, 0.01,""),
    (31, 5,"[12 7 7 1 6]",  "2.78/0.60/0.16/1.22/1.12","2.1698","2.1698","YES", 70, 5645,  8,1263, 0.03,28.66,""),
    (32, 6,"[1 12 9 11 8 10]","1.01/1.05/1.17/2.26/1.11/0.73","3.3193","NaN","N/A",3190,6908,404,  0, 7.05,30.03,"OFF"),
    (33, 6,"[11 5 8 4 3 5]","0.98/2.65/2.47/2.13/2.88/1.27","3.2343","NaN","N/A",509,5986, 68,  0, 0.36,30.00,"OFF"),
    (34, 3,"[2 4 2]",       "0.28/1.81/1.09",          "1.4372","1.4372","YES", 36,   92,  6,  19, 0.01, 0.02,""),
    (35, 5,"[3 9 3 4 9]",   "1.68/1.00/1.63/2.08/2.74","0.7352","0.7352","YES", 22,  114,  2,  24, 0.00, 0.02,""),
    (36, 5,"[3 9 10 3 10]", "1.39/2.34/0.71/1.00/2.86","0.3883","0.3883","YES", 26,   42,  2,   4, 0.01, 0.01,""),
    (37, 5,"[10 9 3 6 10]", "2.03/0.02/0.26/1.04/2.83","0.0000","0.0000","YES", 15,   15,  1,   1, 0.00, 0.00,""),
    (38, 4,"[4 5 3 6]",     "0.65/2.54/1.37/0.84",     "0.6001","0.6001","YES", 20,   46,  2,   7, 0.00, 0.01,""),
    (39, 6,"[4 11 1 9 6 6]","0.11/0.12/1.00/2.84/1.85/1.11","1.1628","1.1628","YES", 59,5376,  6,   1, 0.04,30.01,"OFF"),
    (40, 5,"[3 2 5 11 7]",  "0.89/2.85/2.45/0.97/2.92","0.3720","0.3720","YES", 55,  691,  6, 165, 0.02, 0.56,""),
    (41, 6,"[5 8 5 4 1 4]", "0.81/1.20/0.55/2.86/0.31/1.88","NaN","NaN","N/A",5578,6952,  0,   0,30.01,30.00,"ON OFF"),
    (42, 4,"[6 5 11 4]",    "0.06/2.75/2.59/0.83",     "0.6263","0.6263","YES", 38,   70,  4,  10, 0.03, 0.03,""),
    (43, 5,"[2 2 11 5 8]",  "2.83/0.74/0.04/0.07/2.13","0.6959","0.6959","YES", 30,   77,  3,  13, 0.02, 0.05,""),
    (44, 6,"[6 5 7 11 8 3]","2.32/0.40/0.50/1.84/0.72/2.11","3.4409","NaN","N/A",521,5442, 63,  0, 1.06,30.00,"OFF"),
    (45, 4,"[4 12 1 8]",    "0.12/2.28/0.69/0.27",     "1.3791","1.3791","YES", 54,  100, 10,  27, 0.06, 0.06,""),
    (46, 5,"[9 9 1 4 6]",   "0.86/2.43/0.39/1.84/2.96","0.3180","0.3180","YES", 22,   69,  2,  12, 0.02, 0.03,""),
    (47, 6,"[3 1 12 11 12 5]","2.23/0.64/1.18/2.55/0.38/2.68","0.0000","0.0000","YES", 17, 17,  1,   1, 0.00, 0.00,""),
    (48, 4,"[6 4 12 7]",    "2.41/2.57/2.77/0.91",     "0.6245","0.6245","YES", 19,   52,  2,  10, 0.01, 0.01,""),
    (49, 4,"[8 6 12 5]",    "1.43/1.85/1.21/2.98",     "0.0000","0.0000","YES", 11,   11,  1,   1, 0.00, 0.00,""),
    (50, 3,"[3 4 2]",       "0.85/2.34/1.57",          "0.5545","0.5545","YES", 25,   44,  4,   8, 0.01, 0.01,""),
]

HDR    = ["#","N","Routes","d1 (s)","Cost ON","Cost OFF","Match","Nodes ON","Nodes OFF","Lvs ON","Lvs OFF","t_ON(s)","t_OFF(s)","Tout"]
COL_W  = [0.30, 0.25, 1.20, 2.10, 0.72, 0.72, 0.55, 0.72, 0.75, 0.55, 0.55, 0.65, 0.72, 0.40]

def table_slide(title, subtitle, data):
    sl = prs.slides.add_slide(BLANK)
    bg(sl); title_bar(sl, title, subtitle)
    x0, y0, rh = 0.18, 1.25, 0.455
    xc = x0
    for h, cw in zip(HDR, COL_W):
        rect(sl, xc, y0, cw-0.02, rh-0.04, C_BLUE)
        box(sl, xc+0.02, y0+0.06, cw-0.06, rh-0.12, h,
            font_size=9, bold=True, fg=C_WHITE, align=PP_ALIGN.CENTER)
        xc += cw
    for ri, row in enumerate(data):
        y = y0 + (ri+1)*rh
        bg_row = RGBColor(0xF4,0xF6,0xF8) if ri%2==0 else RGBColor(0xFF,0xFF,0xFF)
        match = row[6]
        if   match == "NO★": bg_row = RGBColor(0xFC,0xE8,0xE6)
        elif match == "N/A":  bg_row = RGBColor(0xFE,0xF5,0xE4)
        xc = x0
        vals = [str(row[0]), str(row[1]), row[2], row[3],
                row[4], row[5], row[6],
                str(row[7]), str(row[8]), str(row[9]), str(row[10]),
                f"{row[11]:.2f}", f"{row[12]:.2f}", row[13]]
        for ci, (v, cw) in enumerate(zip(vals, COL_W)):
            rect(sl, xc, y, cw-0.02, rh-0.04, bg_row)
            fc = C_LGRAY
            if   ci==6 and v=="YES":   fc = C_GREEN
            elif ci==6 and v=="NO★":  fc = C_RED
            elif ci==6 and v=="N/A":   fc = C_YELLOW
            elif ci==13 and v in ("OFF","ON OFF"): fc = C_RED
            box(sl, xc+0.01, y+0.07, cw-0.04, rh-0.14, v,
                font_size=9, fg=fc, align=PP_ALIGN.CENTER)
            xc += cw

# ── Divider slide ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
rect(sl, 0, 2.8, 13.33, 0.08, C_BLUE)
box(sl, 1, 1.5, 11.3, 1.0,
    "Simulation Test 2",
    font_size=40, bold=True, fg=C_BLUE, align=PP_ALIGN.CENTER)
box(sl, 1, 2.6, 11.3, 0.5,
    "50 random trials  ·  N_sys ∈ {3–6}  ·  routes ∈ {1–12}  ·  d1 ∈ U[0, 3 s]  ·  seed = 123",
    font_size=16, fg=C_MGRAY, align=PP_ALIGN.CENTER)

# ── Tables (trials 1-17, 18-34, 35-50) ───────────────────────────────────────
table_slide("Simulation Test 2 — Trials 1–17",
            "50 random cases · WeakRule ON vs OFF · seed=123", rows_all[:17])
table_slide("Simulation Test 2 — Trials 18–34",
            "50 random cases · WeakRule ON vs OFF · seed=123", rows_all[17:34])
table_slide("Simulation Test 2 — Trials 35–50",
            "50 random cases · WeakRule ON vs OFF · seed=123", rows_all[34:])

# ── Summary slide ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "Simulation Test 2 — Summary Statistics",
          "50 trials · seed=123")

stats = [
    ("36 / 50", "Cost identical\n(ON = OFF)", C_GREEN),
    ("2 / 50",  "Cost differs\n(Trials 3 & 23 — ON cheaper)", C_YELLOW),
    ("12 / 50", "N/A — OFF timed out\n(ON completed in 10/12)", C_RED),
]
for i, (val, label, col) in enumerate(stats):
    x = 0.7 + i * 4.1
    rect(sl, x, 1.3, 3.6, 1.7, col)
    box(sl, x, 1.4, 3.6, 0.9, val, font_size=44, bold=True,
        fg=C_WHITE, align=PP_ALIGN.CENTER)
    box(sl, x, 2.2, 3.6, 0.7, label, font_size=13,
        fg=C_WHITE, align=PP_ALIGN.CENTER)

box(sl, 0.5, 3.2, 12.3, 0.4,
    "Key observation: WeakRule ON is NEVER worse than OFF — in both cost-differ cases, ON finds the better solution",
    font_size=14, bold=True, fg=C_YELLOW, align=PP_ALIGN.CENTER)

perf = [
    ("Avg time ON",  "2.30 s",  C_GREEN),
    ("Avg time OFF", "9.07 s",  C_RED),
    ("Speedup (avg)","~4×",     C_BLUE),
    ("ON timeouts",  "2 / 50",  C_YELLOW),
    ("OFF timeouts", "12 / 50", C_RED),
]
for i, (k, v, col) in enumerate(perf):
    x = 0.4 + i * 2.5
    rect(sl, x, 3.8, 2.3, 0.5, RGBColor(0xF4,0xF6,0xF8))
    rect(sl, x, 3.8, 0.08, 0.5, col)
    box(sl, x+0.15, 3.82, 2.1, 0.25, k, font_size=10, fg=C_MGRAY)
    box(sl, x+0.15, 4.08, 2.1, 0.25, v, font_size=13, bold=True, fg=col)

bullets = [
    "• Trial 3:  routes=[8 9 4 5 3]  ON=2.7792 < OFF=2.9938  — same V_temp/x mechanism as Test 1 Trial 13",
    "• Trial 23: routes=[8 1 6 12 5] ON=1.7178 < OFF=1.7974  — ON again finds strictly better solution",
    "• Trials 6 & 41: both ON and OFF timeout (N=6, high contention) — neither finds complete path",
    "• N=6 cases dominate timeouts; N≤4 cases all complete rapidly with matching costs",
    "• Combined across Test 1+2 (70 trials): 51 YES · 3 NO (ON always better) · 16 N/A",
]
for i, b in enumerate(bullets):
    box(sl, 0.5, 4.55 + i*0.47, 12.3, 0.42, b, font_size=12, fg=C_LGRAY)

# ── save ─────────────────────────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f"Appended 5 slides to: {PPTX_PATH}")
print(f"Total slides now: {len(prs.slides)}")
