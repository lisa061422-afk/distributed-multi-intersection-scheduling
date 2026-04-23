"""
make_algo_ppt.py — generates Algorithm_Modifications.pptx
Four core modifications of WeakRule ON, paper-style, English.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BLUE   = RGBColor(0x1F, 0x6B, 0xB0)
C_GREEN  = RGBColor(0x1E, 0x8A, 0x44)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_ORANGE = RGBColor(0xD3, 0x57, 0x00)
C_PURPLE = RGBColor(0x6C, 0x35, 0x9E)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x2C, 0x3E, 0x50)
C_GRAY   = RGBColor(0x55, 0x6B, 0x82)
C_LGRAY  = RGBColor(0xF4, 0xF6, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────────
def bg(slide, color=C_WHITE):
    sh = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()

def rect(slide, l, t, w, h, fill, line_color=None, line_w=Pt(0)):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color: sh.line.color.rgb = line_color; sh.line.width = line_w
    else: sh.line.fill.background()
    return sh

def box(slide, l, t, w, h, text, size=13, bold=False,
        fg=C_DARK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = fg
    return txb

def multiline(slide, l, t, w, lines, size=12, fg=C_DARK, indent=False):
    """lines: list of (text, bold, color_override_or_None)"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(6))
    tf  = txb.text_frame; tf.word_wrap = True
    first = True
    for text, bold, color in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color if color else fg

def title_bar(slide, title, subtitle=None, bar_color=C_BLUE, tag=None, tag_color=None):
    rect(slide, 0, 0, 13.33, 1.05, bar_color)
    box(slide, 0.35, 0.08, 11.5, 0.6, title,
        size=28, bold=True, fg=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        box(slide, 0.35, 0.65, 10, 0.38, subtitle,
            size=13, fg=RGBColor(0xD0,0xE8,0xFF), align=PP_ALIGN.LEFT)
    if tag:
        tc = tag_color or bar_color
        rect(slide, 11.8, 0.15, 1.3, 0.45, RGBColor(0xFF,0xFF,0xFF))
        box(slide, 11.8, 0.18, 1.3, 0.38, tag,
            size=13, bold=True, fg=tc, align=PP_ALIGN.CENTER)

def accent_block(slide, l, t, w, h, title, body_lines, accent_color,
                 title_size=14, body_size=12):
    rect(slide, l, t, 0.08, h, accent_color)
    rect(slide, l+0.08, t, w-0.08, h, C_LGRAY)
    box(slide, l+0.18, t+0.08, w-0.28, 0.32,
        title, size=title_size, bold=True, fg=accent_color)
    y = t + 0.38
    for line in body_lines:
        box(slide, l+0.18, y, w-0.28, 0.3, line, size=body_size, fg=C_DARK)
        y += 0.28

def pill(slide, l, t, text, color):
    rect(slide, l, t, 1.6, 0.32, color)
    box(slide, l, t+0.04, 1.6, 0.25, text,
        size=11, bold=True, fg=C_WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK); bg(sl)
rect(sl, 0, 0, 13.33, 7.5, RGBColor(0xF8,0xF9,0xFA))
rect(sl, 0, 0, 0.18, 7.5, C_BLUE)
rect(sl, 0, 3.2, 13.33, 0.06, C_BLUE)

box(sl, 0.5, 0.8, 12.3, 1.1,
    "WeakRule ON Algorithm", size=42, bold=True, fg=C_BLUE, align=PP_ALIGN.LEFT)
box(sl, 0.5, 1.9, 12.3, 0.6,
    "Four Core Modifications for Optimality and Efficiency",
    size=22, fg=C_GRAY, align=PP_ALIGN.LEFT)
rect(sl, 0.5, 2.65, 3.5, 0.06, C_BLUE)

mods = [
    (C_BLUE,   "Mod 1", "Pairwise Priority Lock"),
    (C_GREEN,  "Mod 2", "x / V_temp Co-Decision"),
    (C_ORANGE, "Mod 3", "reset_since — Stale Record Tracking"),
    (C_PURPLE, "Mod 4", "T_bound Pruning"),
]
for i, (col, tag, desc) in enumerate(mods):
    x = 0.5 + i * 3.2
    rect(sl, x, 3.5, 3.0, 0.9, col)
    box(sl, x+0.12, 3.58, 2.76, 0.35, tag,
        size=16, bold=True, fg=C_WHITE)
    box(sl, x+0.12, 3.9, 2.76, 0.4, desc,
        size=11, fg=RGBColor(0xE0,0xF0,0xFF))

box(sl, 0.5, 6.9, 12.3, 0.4,
    "Single-Intersection Centralized Decision Tree Scheduler  ·  April 2026",
    size=11, fg=C_GRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  Mod 1 — Pairwise Priority Lock
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK); bg(sl)
title_bar(sl,
    "Mod 1 — Pairwise Priority Lock",
    "Preventing redundant branching and cycling in resource contention",
    bar_color=C_BLUE, tag="Mod 1 / 4", tag_color=C_BLUE)

# Problem
rect(sl, 0.35, 1.18, 1.1, 0.38, C_RED)
box(sl, 0.35, 1.21, 1.1, 0.3, "PROBLEM", size=11, bold=True,
    fg=C_WHITE, align=PP_ALIGN.CENTER)
box(sl, 1.55, 1.18, 11.2, 0.38,
    "Each time two systems contest the same resource, a new branch is created → "
    "exponential tree growth and cycling.",
    size=12, fg=C_DARK)

# Key data structure
rect(sl, 0.35, 1.75, 12.5, 0.38, C_BLUE)
box(sl, 0.35, 1.78, 12.5, 0.3,
    "Data Structure:  pair_lock(i, j) = winner  —  N × N integer matrix",
    size=13, bold=True, fg=C_WHITE)

# Semantics table
headers = ["pair_lock(i, j)", "Meaning", "Effect at next contention"]
col_w = [2.5, 4.5, 5.3]; x0 = 0.35; y0 = 2.22; rh = 0.42
xc = x0
for h, cw in zip(headers, col_w):
    rect(sl, xc, y0, cw-0.04, rh-0.04, RGBColor(0xD6,0xE8,0xF8))
    box(sl, xc+0.08, y0+0.08, cw-0.16, 0.28, h, size=11, bold=True, fg=C_BLUE)
    xc += cw
rows_t = [
    ("0", "Never competed", "Generate two branches (enumerate both priorities)"),
    ("= j", "j won over i", "j auto-wins — no new branch created"),
    ("= i", "i won over j", "i auto-wins — no new branch created"),
]
for ri, (a,b,c) in enumerate(rows_t):
    y = y0 + (ri+1)*rh
    bg_r = C_LGRAY if ri%2==0 else C_WHITE
    xc = x0
    for vi, (v, cw) in enumerate(zip([a,b,c], col_w)):
        rect(sl, xc, y, cw-0.04, rh-0.04, bg_r)
        fc = C_GREEN if vi==2 and "auto" in v else C_DARK
        box(sl, xc+0.08, y+0.1, cw-0.16, 0.28, v, size=11, fg=fc)
        xc += cw

# Auto-win condition
rect(sl, 0.35, 3.58, 12.5, 0.35, C_GREEN)
box(sl, 0.35, 3.61, 12.5, 0.28,
    "Auto-win condition:  pair_lock(c, o) = c  for ALL contenders o  AND  ra[s_c, c] > 0  (actively executing)",
    size=12, bold=True, fg=C_WHITE)

# Two key points
pts = [
    (C_BLUE,  "Completeness guarantee",
     "All priority combinations are enumerated at the first competition "
     "(pair_lock = 0 → two branches). Subsequent lock enforcement only prevents "
     "re-branching on the same path; sibling branches remain independent."),
    (C_ORANGE,"Key design decision",
     "Locks are NEVER cleared on reset. Clearing was the root cause of cycling "
     "in the prior version: cleared locks → repeated competition on the same path → infinite branches."),
    (C_PURPLE,"Third-system entry",
     "If system n3 enters with pair_lock(n1, n3) = 0, auto-win does not apply "
     "→ new branches are automatically created to cover all orderings involving n3."),
]
for i, (col, title, body) in enumerate(pts):
    y = 4.08 + i * 0.98
    rect(sl, 0.35, y, 0.1, 0.75, col)
    box(sl, 0.55, y+0.04, 3.2, 0.28, title, size=12, bold=True, fg=col)
    box(sl, 0.55, y+0.32, 12.1, 0.5, body, size=11, fg=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3  Mod 2 — x / V_temp Co-Decision
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK); bg(sl)
title_bar(sl,
    "Mod 2 — x / V_temp Co-Decision  (space_variants)",
    "Allowing reset sub-tasks to displace V_temp occupants → strictly better solutions",
    bar_color=C_GREEN, tag="Mod 2 / 4", tag_color=C_GREEN)

rect(sl, 0.35, 1.18, 1.1, 0.38, C_RED)
box(sl, 0.35, 1.21, 1.1, 0.3, "PROBLEM", size=11, bold=True,
    fg=C_WHITE, align=PP_ALIGN.CENTER)
box(sl, 1.55, 1.18, 11.2, 0.38,
    "Original method: reset sub-task schedule x depends unilaterally on current V_temp occupant "
    "→ cannot reach solutions requiring V_temp reassignment (suboptimal in some cases).",
    size=12, fg=C_DARK)

rect(sl, 0.35, 1.72, 12.5, 0.35, C_GREEN)
box(sl, 0.35, 1.75, 12.5, 0.28,
    "Modification:  space_variants  generates multiple branches when reset system n needs space m (occupied by np)",
    size=13, bold=True, fg=C_WHITE)

# Branch table
headers2 = ["pair_lock(n, np)", "Branch type", "Action", "Result"]
col_w2 = [2.3, 2.1, 4.5, 3.3]; x0=0.35; y0=2.18; rh=0.44
xc = x0
for h, cw in zip(headers2, col_w2):
    rect(sl, xc, y0, cw-0.04, rh-0.04, RGBColor(0xD0,0xF0,0xDC))
    box(sl, xc+0.08, y0+0.09, cw-0.16, 0.28, h, size=11, bold=True, fg=C_GREEN)
    xc += cw
rows2 = [
    ("0  (never competed)", "wait + displace", "Both variants → 2 new branches", "Both orderings explored"),
    ("= np  (np won)", "wait only", "np keeps space m", "n waits for np to finish"),
    ("= n   (n won)", "displace only", "np evicted from V_temp, fully reset", "n immediately re-enters m"),
]
for ri, row in enumerate(rows2):
    y = y0 + (ri+1)*rh
    bg_r = C_LGRAY if ri%2==0 else C_WHITE
    xc = x0
    for ci, (v, cw) in enumerate(zip(row, col_w2)):
        rect(sl, xc, y, cw-0.04, rh-0.04, bg_r)
        fc = C_GREEN if ci==1 else C_DARK
        box(sl, xc+0.08, y+0.1, cw-0.16, 0.3, v, size=10, fg=fc)
        xc += cw

pts2 = [
    (C_GREEN,  "Displace branch",
     "np is removed from V_temp and fully reset (ra ← C, x ← {}). "
     "System n's sub-task start time is aligned to the moment space m becomes available, "
     "eliminating idle gaps. This is the mechanism that allows ON to find strictly cheaper "
     "solutions than OFF (verified in Trial 13: g_ON = 1.58 vs g_OFF = 1.72)."),
    (C_BLUE,   "Scope constraint",
     "Only the pending V_temp requests of other systems are modified. "
     "V_valid assignments already committed by traverse_columns in the main branching loop "
     "remain untouched — the co-decision applies only to V_temp, not V_valid."),
    (C_ORANGE, "Why OFF cannot find these solutions",
     "OFF treats V_temp as fixed: x must wait for the current occupant np. "
     "This locks priority to np, preventing the displace branch from being explored."),
]
for i, (col, title, body) in enumerate(pts2):
    y = 3.65 + i * 1.02
    rect(sl, 0.35, y, 0.1, 0.78, col)
    box(sl, 0.55, y+0.04, 3.0, 0.28, title, size=12, bold=True, fg=col)
    box(sl, 0.55, y+0.32, 12.1, 0.55, body, size=11, fg=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4  Mod 3 — reset_since
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK); bg(sl)
title_bar(sl,
    "Mod 3 — reset_since: Stale Occupation Record Tracking",
    "Invalidating historical V entries after a reset without modifying shared parent nodes",
    bar_color=C_ORANGE, tag="Mod 3 / 4", tag_color=C_ORANGE)

rect(sl, 0.35, 1.18, 1.1, 0.38, C_RED)
box(sl, 0.35, 1.21, 1.1, 0.3, "PROBLEM", size=11, bold=True,
    fg=C_WHITE, align=PP_ALIGN.CENTER)
box(sl, 1.55, 1.18, 11.2, 0.38,
    "After system n is reset, its historical V occupation records are invalid for the current branch. "
    "V cannot be directly modified — it is shared across all ancestor nodes.",
    size=12, fg=C_DARK)

rect(sl, 0.35, 1.72, 12.5, 0.35, C_ORANGE)
box(sl, 0.35, 1.75, 12.5, 0.28,
    "Solution:  Each node stores  reset_since[n]  —  a 1×N vector recording the reset timestamp for each system",
    size=13, bold=True, fg=C_WHITE)

# Mechanism
blocks = [
    ("On interrupt at time tw",
     ["Set  reset_since[n] = tw",
      "Clear stale reservations:  x{n}{k} = {}",
      "V is NOT modified  (ancestor nodes share the same V matrix)"]),
    ("During resource check  (check_resc_occupation)",
     ["V entries for system n with timestamp < reset_since[n] are ignored",
      "Only records created after the last reset are considered valid"]),
]
for i, (title, lines) in enumerate(blocks):
    y = 2.22 + i * 1.35
    rect(sl, 0.35, y, 0.08, 1.1, C_ORANGE)
    rect(sl, 0.43, y, 12.42, 1.1, C_LGRAY)
    box(sl, 0.55, y+0.08, 12.1, 0.3, title, size=13, bold=True, fg=C_ORANGE)
    for j, line in enumerate(lines):
        box(sl, 0.65, y+0.42+j*0.3, 12.0, 0.28,
            "• " + line, size=11, fg=C_DARK)

# Schedule reconstruction
rect(sl, 0.35, 4.95, 12.5, 0.35, RGBColor(0x5D,0x4A,0x1A))
box(sl, 0.35, 4.98, 12.5, 0.28,
    "Final schedule reconstruction from a complete leaf node",
    size=13, bold=True, fg=C_WHITE)

recon = [
    ("gamma[n]",  C_GREEN,
     "Authoritative occupation intervals written when each sub-task completes. "
     "Primary source for the scheduling timeline."),
    ("x{n}{k}",   C_BLUE,
     "Reservation records generated during resets: {t_start, t_end, sub_index, space}. "
     "Covers sub-tasks that were rescheduled but not yet completed at the leaf."),
    ("Filter",     C_ORANGE,
     "Merge gamma and x; discard any entry for system n with timestamp < reset_since[n]. "
     "The result is the valid, non-overlapping schedule for each system."),
]
for i, (label, col, body) in enumerate(recon):
    y = 5.42 + i * 0.62
    rect(sl, 0.35, y, 1.2, 0.42, col)
    box(sl, 0.35, y+0.07, 1.2, 0.28, label, size=12, bold=True,
        fg=C_WHITE, align=PP_ALIGN.CENTER)
    box(sl, 1.65, y+0.07, 11.1, 0.42, body, size=11, fg=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5  Mod 4 — T_bound
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK); bg(sl)
title_bar(sl,
    "Mod 4 — T_bound: Serial Worst-Case Pruning",
    "Removing open nodes that cannot improve over a naive serial baseline",
    bar_color=C_PURPLE, tag="Mod 4 / 4", tag_color=C_PURPLE)

rect(sl, 0.35, 1.18, 1.1, 0.38, C_RED)
box(sl, 0.35, 1.21, 1.1, 0.3, "PROBLEM", size=11, bold=True,
    fg=C_WHITE, align=PP_ALIGN.CENTER)
box(sl, 1.55, 1.18, 11.2, 0.38,
    "The open list accumulates nodes whose schedules are provably worse than executing "
    "all tasks serially, wasting search time on paths that can never be competitive.",
    size=12, fg=C_DARK)

# Formula block
rect(sl, 0.35, 1.72, 12.5, 1.55, RGBColor(0xF0,0xEB,0xF8))
rect(sl, 0.35, 1.72, 0.08, 1.55, C_PURPLE)
box(sl, 0.55, 1.78, 12.1, 0.35,
    "Worst-case serial completion time:", size=13, bold=True, fg=C_PURPLE)
box(sl, 0.55, 2.15, 12.1, 0.45,
    "T_bound  =  max_n ( α̃_n )  +  Σ_{n=1}^{N} Σ_{s=1}^{S} C_{s,n}",
    size=18, bold=True, fg=C_DARK, align=PP_ALIGN.CENTER)
box(sl, 0.55, 2.62, 12.1, 0.55,
    "where  α̃_n = d1_n + (R/2 − W/2) / v_max  is the earliest arrival time of system n at the merge zone",
    size=11, fg=C_GRAY, align=PP_ALIGN.CENTER)

# Pruning rule
rect(sl, 0.35, 3.38, 12.5, 0.45, C_PURPLE)
box(sl, 0.35, 3.43, 12.5, 0.35,
    "Pruning rule:  remove any node from OPEN where  tw > T_bound",
    size=14, bold=True, fg=C_WHITE, align=PP_ALIGN.CENTER)

# Interpretation
pts3 = [
    (C_PURPLE, "Interpretation",
     "T_bound = time if the last-arriving system waits for ALL other systems' sub-tasks "
     "to complete serially before starting. Any valid cooperative schedule must finish "
     "earlier than this bound."),
    (C_BLUE,   "Validity",
     "tw in the search tree represents the current decision point time. "
     "A node with tw > T_bound cannot possibly produce a schedule better than the serial "
     "baseline, so all descendants are safely pruned."),
    (C_GREEN,  "Application scope",
     "Applied to both WeakRule ON and OFF in the current implementation. "
     "For WeakRule OFF, tw is monotonically non-decreasing, making the bound especially tight. "
     "For WeakRule ON, displacement may cause non-monotone tw, but the bound remains valid "
     "as a sufficient (conservative) pruning condition."),
]
for i, (col, title, body) in enumerate(pts3):
    y = 3.98 + i * 1.05
    rect(sl, 0.35, y, 0.1, 0.82, col)
    box(sl, 0.55, y+0.04, 3.0, 0.28, title, size=12, bold=True, fg=col)
    box(sl, 0.55, y+0.34, 12.1, 0.55, body, size=11, fg=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6  Summary
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK); bg(sl)
title_bar(sl, "Summary — Four Core Modifications",
          "WeakRule ON vs. original method", bar_color=C_DARK)

cols_h = ["Modification", "Problem Addressed", "Mechanism", "Guarantee"]
col_w3 = [2.2, 3.0, 4.3, 3.5]; x0=0.25; y0=1.18; rh=0.5
xc = x0
for h, cw in zip(cols_h, col_w3):
    rect(sl, xc, y0, cw-0.04, rh-0.04, C_DARK)
    box(sl, xc+0.08, y0+0.1, cw-0.16, 0.32, h,
        size=12, bold=True, fg=C_WHITE)
    xc += cw

rows_s = [
    (C_BLUE,   "Mod 1\nPairwise Priority Lock",
     "Tree explosion & cycling from repeated contention",
     "N×N pair_lock records pairwise winners; auto-win skips branching if candidate beats all via established locks",
     "All orderings explored once at first encounter; no redundant re-branching"),
    (C_GREEN,  "Mod 2\nx / V_temp Co-Decision",
     "Reset schedule x constrained by fixed V_temp → suboptimal solutions",
     "space_variants generates wait and displace branches; displace evicts V_temp occupant and aligns reset timing",
     "ON can find solutions strictly better than OFF (e.g. Trial 13: 1.58 vs 1.72)"),
    (C_ORANGE, "Mod 3\nreset_since",
     "Stale V records persist after reset; V shared across nodes cannot be modified",
     "reset_since[n]=tw marks reset time; V entries before tw ignored; x cleared; schedule rebuilt from gamma+x",
     "Correct schedule reconstruction without corrupting shared ancestor state"),
    (C_PURPLE, "Mod 4\nT_bound Pruning",
     "Open list polluted with nodes worse than serial execution",
     "T_bound = max(α̃_n) + ΣC; nodes with tw > T_bound removed from OPEN",
     "Search focused on competitive regions; provably safe pruning"),
]
for ri, (col, mod, prob, mech, guar) in enumerate(rows_s):
    y = y0 + (ri+1)*rh
    bg_r = C_LGRAY if ri%2==0 else C_WHITE
    xc = x0
    for ci, (v, cw) in enumerate(zip([mod, prob, mech, guar], col_w3)):
        rect(sl, xc, y, cw-0.04, rh-0.04, bg_r)
        if ci == 0:
            rect(sl, xc, y, 0.07, rh-0.04, col)
        fc = col if ci == 0 else C_DARK
        box(sl, xc+0.12, y+0.06, cw-0.2, rh-0.14, v,
            size=10, bold=(ci==0), fg=fc)
        xc += cw

box(sl, 0.25, 6.3, 12.8, 0.4,
    "Result across 90 random trials (Test 1–3):  51 cost-match · 3 cases ON strictly better · 0 cases ON suboptimal",
    size=13, bold=True, fg=C_BLUE, align=PP_ALIGN.CENTER)
box(sl, 0.25, 6.75, 12.8, 0.35,
    "WeakRule ON is never worse than OFF, and achieves up to 16× speedup on N=3–6 systems.",
    size=12, fg=C_GRAY, align=PP_ALIGN.CENTER)

# ── save ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\robin\OneDrive\Documents\Github_file\Astar_0619_resettingrule_dif_x\Algorithm_Modifications.pptx"
prs.save(out)
print("Saved:", out)
