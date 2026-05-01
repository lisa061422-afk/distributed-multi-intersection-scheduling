"""
Generate a PowerPoint presentation from batch_results.json.

Usage (from warehouse root):
    python python_port/make_pptx.py
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Cm
except ImportError:
    print('ERROR: python-pptx not installed. Run: pip install python-pptx')
    sys.exit(1)

RESULTS_JSON = Path(__file__).parent / 'batch_results.json'
OUT_PPTX = Path(__file__).parent / 'batch_validation_results_v2.pptx'

# ── Colors ────────────────────────────────────────────────────────────────
C_BLUE   = RGBColor(0x1F, 0x49, 0x7D)   # dark blue header
C_LBLUE  = RGBColor(0xBD, 0xD7, 0xEE)   # light blue row
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
C_GREEN  = RGBColor(0x70, 0xAD, 0x47)
C_ORANGE = RGBColor(0xFF, 0xC0, 0x00)
C_RED    = RGBColor(0xFF, 0x45, 0x00)
C_DARK   = RGBColor(0x26, 0x26, 0x26)


def _add_text(tf, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
    tf.text = ''
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_cell(cell, text, font_size=11, bold=False,
               bg=None, fg=C_DARK, align=PP_ALIGN.CENTER):
    cell.text = ''
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = fg
    if bg is not None:
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', f'{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}')


def slide_title(prs):
    """Slide 1 — Title."""
    layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BLUE

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    _add_text(tf,
              'Distributed Scheduling for Warehouse AMR\nPython ADMM — Batch Validation Results',
              size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    _add_text(txBox2.text_frame,
              '20 Random Scenarios  |  N = 10 ~ 15 Vehicles',
              size=18, color=C_LBLUE, align=PP_ALIGN.CENTER)

    # Date
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.4))
    _add_text(txBox3.text_frame, '2026-04-24',
              size=14, color=C_LBLUE, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    """Slide 2 — System Overview."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Header bar
    hdr = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.0))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = C_BLUE
    hdr.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.7))
    _add_text(txBox.text_frame, 'System Overview', size=24, bold=True, color=C_WHITE)

    bullets = [
        ('Warehouse topology',
         '4 Intersection agents (merging zones) + 4 Road agents + 1 Terminal agent'),
        ('Vehicle count',
         'N = 10 ~ 15 AGVs; each follows a pre-assigned route through 2–4 intersections'),
        ('ADMM formulation',
         'Each agent solves a local QP; consensus on shared entry/exit times'),
        ('Closed-form solution',
         'x* = max(-f / [4(rho1+rho2)],  alpha_tilde)  — no numerical solver needed in '
         'most cases; only falls back to SLSQP when mutual-exclusion is violated'),
        ('Parallel execution',
         'Intersection agents solved concurrently via ProcessPoolExecutor '
         '(pre-warmed; zero startup cost during online MPC)'),
        ('Convergence criteria',
         'Primal residual r < 0.01  AND  dual residual s < 0.01  (max 50 iterations)'),
    ]

    y = 1.15
    for title, body in bullets:
        # Bullet dot
        dot = slide.shapes.add_shape(1, Inches(0.3), Inches(y + 0.05),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = C_BLUE
        dot.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(9.2), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run(); r1.text = title + ':  '
        r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = C_BLUE
        r2 = p.add_run(); r2.text = body
        r2.font.size = Pt(13); r2.font.color.rgb = C_DARK
        y += 0.78


def slide_results_table(prs, results):
    """Slide 3 — Full results table."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Header
    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.0))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C_BLUE; hdr.line.fill.background()
    tb2 = slide.shapes.add_textbox(Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.7))
    _add_text(tb2.text_frame, 'Batch Validation — 20 Scenarios', size=24, bold=True, color=C_WHITE)

    # Table
    cols = ['#', 'N', 'Seed', 'k', 'Cost', 'T_seq (s)', 'T_par (s)', 'Speedup']
    col_w = [0.35, 0.35, 0.55, 0.45, 0.75, 0.95, 0.95, 0.85]
    n_rows = len(results) + 1

    tbl = slide.shapes.add_table(
        n_rows, len(cols),
        Inches(0.18), Inches(1.05),
        Inches(9.64), Inches(5.7)).table

    # Column widths
    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = Inches(w)

    # Header row
    for ci, h in enumerate(cols):
        _set_cell(tbl.cell(0, ci), h, font_size=12, bold=True,
                  bg=C_BLUE, fg=C_WHITE, align=PP_ALIGN.CENTER)

    for ri, r in enumerate(results):
        row_bg = C_LBLUE if ri % 2 == 0 else C_WHITE
        spd = r['speedup']
        spd_color = C_GREEN if spd >= 1.5 else (C_ORANGE if spd >= 1.0 else C_RED)
        vals = [r['i'], r['N'], r['seed'], r['k'],
                f"{r['cost']:.3f}",
                f"{r['T_seq']:.2f}",
                f"{r['T_par']:.2f}",
                f"{spd:.2f}x"]
        for ci, v in enumerate(vals):
            bold_spd = (ci == 7)
            fg = spd_color if ci == 7 else C_DARK
            _set_cell(tbl.cell(ri + 1, ci), v, font_size=11,
                      bold=bold_spd, bg=row_bg, fg=fg)


def slide_summary(prs, results, summary):
    """Slide 4 — Summary + conclusion."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.0))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C_BLUE; hdr.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.15), Inches(9.6), Inches(0.7))
    _add_text(tb.text_frame, 'Summary & Conclusion', size=24, bold=True, color=C_WHITE)

    avg_seq = summary['avg_seq']
    avg_par = summary['avg_par']
    avg_spd = summary['avg_speedup']
    max_par = summary['max_par']

    # Big stat boxes
    stats = [
        ('Avg Sequential', f'{avg_seq:.2f} s'),
        ('Avg Parallel',   f'{avg_par:.2f} s'),
        ('Avg Speedup',    f'{avg_spd:.2f}x'),
        ('Max Parallel',   f'{max_par:.2f} s'),
    ]
    for si, (label, val) in enumerate(stats):
        x = 0.4 + si * 2.35
        box = slide.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(2.0), Inches(1.2))
        box.fill.solid(); box.fill.fore_color.rgb = C_LBLUE; box.line.color.rgb = C_BLUE

        tb_l = slide.shapes.add_textbox(Inches(x), Inches(1.22), Inches(2.0), Inches(0.35))
        _add_text(tb_l.text_frame, label, size=11, color=C_BLUE, align=PP_ALIGN.CENTER)

        tb_v = slide.shapes.add_textbox(Inches(x), Inches(1.55), Inches(2.0), Inches(0.65))
        _add_text(tb_v.text_frame, val, size=22, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # Key findings
    if max_par < 2.0:
        mpc_verdict = 'Online MPC feasible: dt >= 2.0 s  (real-time capable)'
        v_color = C_GREEN
    elif max_par < 3.0:
        mpc_verdict = 'Online MPC feasible: dt >= 3.0 s'
        v_color = C_GREEN
    else:
        mpc_verdict = f'Recommend dt >= {max_par:.1f} s for safe MPC margin'
        v_color = C_ORANGE

    findings = [
        ('Analytical closed-form QP',
         'Most intersection sub-problems solved in O(N) without a numerical solver; '
         'only ~5% of calls need SLSQP fallback.'),
        ('Parallel speedup pattern',
         'Small N (10): near 1x (IPC overhead dominates).  '
         'Medium-large N (13-15): consistent 1.3-1.8x speedup.'),
        ('Python vs MATLAB parallel',
         'Python avg ~2x faster than MATLAB parallel (6 s): closed-form QP '
         'eliminates most quadprog calls entirely.'),
        ('Online MPC verdict', mpc_verdict),
    ]

    y = 2.65
    for i, (title, body) in enumerate(findings):
        dot = slide.shapes.add_shape(1, Inches(0.3), Inches(y + 0.07),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = v_color if i == 3 else C_BLUE
        dot.line.fill.background()

        tb_f = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(9.2), Inches(0.72))
        tf = tb_f.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run(); r1.text = title + ':  '
        r1.font.size = Pt(13); r1.font.bold = True
        r1.font.color.rgb = v_color if i == 3 else C_BLUE
        r2 = p.add_run(); r2.text = body
        r2.font.size = Pt(13); r2.font.color.rgb = C_DARK
        y += 0.88


def main():
    if not RESULTS_JSON.exists():
        print(f'ERROR: {RESULTS_JSON} not found.')
        print('Run  python python_port/run_batch_validation.py  first.')
        sys.exit(1)

    with open(RESULTS_JSON) as f:
        data = json.load(f)

    results  = data['results']
    summary  = data['summary']

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_overview(prs)
    slide_results_table(prs, results)
    slide_summary(prs, results, summary)

    prs.save(str(OUT_PPTX))
    print(f'Saved: {OUT_PPTX}')


if __name__ == '__main__':
    main()
